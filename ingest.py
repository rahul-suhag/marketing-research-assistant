"""
ingest.py — Build a FAISS vector index from PDFs, PPTX slides, and text files.

Reads from:  data/slides/  and  data/marketing_books/
Supports:    .pdf, .pptx, .txt, .md

Usage:
    python ingest.py
"""

from pathlib import Path

from langchain_community.document_loaders import PyPDFLoader, TextLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.schema import Document
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS


# ── PPTX extraction helpers ─────────────────────────────────────────────────

def _extract_text_from_shape(shape) -> list:
    """Recursively extract text from a PPTX shape, including groups and tables."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    texts = []
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            t = para.text.strip()
            if t:
                texts.append(t)
    if shape.has_table:
        for row in shape.table.rows:
            row_text = " | ".join(
                cell.text.strip() for cell in row.cells if cell.text.strip()
            )
            if row_text:
                texts.append(row_text)
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            texts.extend(_extract_text_from_shape(child))
    return texts


def load_pptx(pptx_path: Path) -> list:
    """Extract text from a .pptx file, one Document per slide."""
    from pptx import Presentation

    prs = Presentation(str(pptx_path))
    session_name = pptx_path.stem.replace("_", " ").replace("2026", "").strip()

    docs = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            texts.extend(_extract_text_from_shape(shape))
        if texts:
            header = f"[{session_name} — Slide {slide_num}]"
            content = header + "\n" + "\n".join(texts)
            docs.append(Document(
                page_content=content,
                metadata={
                    "source": pptx_path.name,
                    "page": slide_num - 1,
                    "source_category": "slides",
                },
            ))
    return docs

# ── Configuration ────────────────────────────────────────────────────────────

SLIDES_DIR = Path("data/slides")
BOOKS_DIR = Path("data/marketing_books")
INDEX_DIR = Path("faiss_index")

CHUNK_SIZE = 1000
CHUNK_OVERLAP = 200
EMBEDDING_MODEL = "sentence-transformers/all-MiniLM-L6-v2"


def load_documents_from_directory(directory: Path) -> list:
    """Load all PDFs and text files from a directory."""
    documents = []
    if not directory.exists():
        print(f"  Warning: Directory {directory} does not exist — skipping.")
        return documents

    source_tag = "slides" if "slides" in str(directory) else "marketing_books"

    # Load PDFs
    for pdf_path in sorted(directory.glob("*.pdf")):
        print(f"  Loading PDF: {pdf_path.name}")
        loader = PyPDFLoader(str(pdf_path))
        docs = loader.load()
        for doc in docs:
            doc.metadata["source_category"] = source_tag
        documents.extend(docs)

    # Load PPTX files
    for pptx_path in sorted(directory.glob("*.pptx")):
        print(f"  Loading PPTX: {pptx_path.name}")
        docs = load_pptx(pptx_path)
        documents.extend(docs)

    # Load text files (.txt and .md)
    for ext in ("*.txt", "*.md"):
        for txt_path in sorted(directory.glob(ext)):
            print(f"  Loading text: {txt_path.name}")
            loader = TextLoader(str(txt_path), encoding="utf-8")
            docs = loader.load()
            for doc in docs:
                doc.metadata["source_category"] = source_tag
                doc.metadata["source"] = txt_path.name
            documents.extend(docs)

    if not documents:
        print(f"  Warning: No files found in {directory}.")

    return documents


def main():
    print("=" * 60)
    print("FAISS Index Builder — Marketing Research RAG")
    print("=" * 60)

    # ── 1. Load documents ────────────────────────────────────────────────
    print("\n[1/4] Loading documents …")
    all_docs = []
    all_docs.extend(load_documents_from_directory(SLIDES_DIR))
    all_docs.extend(load_documents_from_directory(BOOKS_DIR))

    if not all_docs:
        print("\nNo documents found. Place PDFs or text files in data/slides/ and/or data/marketing_books/.")
        return

    print(f"\n  Total documents loaded: {len(all_docs)}")

    # ── 2. Split into chunks ─────────────────────────────────────────────
    print("\n[2/4] Splitting into chunks …")
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=CHUNK_SIZE,
        chunk_overlap=CHUNK_OVERLAP,
        separators=["\n\n", "\n", ". ", " ", ""],
    )
    chunks = splitter.split_documents(all_docs)
    print(f"  Total chunks created: {len(chunks)}")

    # ── 3. Create embeddings ─────────────────────────────────────────────
    print(f"\n[3/4] Loading embedding model ({EMBEDDING_MODEL}) …")
    embeddings = HuggingFaceEmbeddings(
        model_name=EMBEDDING_MODEL,
        model_kwargs={"device": "cpu"},
        encode_kwargs={"normalize_embeddings": True},
    )

    # ── 4. Build and save FAISS index ────────────────────────────────────
    print("\n[4/4] Building FAISS index …")
    vectorstore = FAISS.from_documents(chunks, embeddings)

    INDEX_DIR.mkdir(exist_ok=True)
    vectorstore.save_local(str(INDEX_DIR))
    print(f"\n  FAISS index saved to {INDEX_DIR}/")
    print(f"  {len(chunks)} chunks indexed.")
    print("  You can now run:  streamlit run app.py")


if __name__ == "__main__":
    main()
